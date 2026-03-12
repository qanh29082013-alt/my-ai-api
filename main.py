"""
AI API - Powered by Claude
Tính năng: Chatbot đa người dùng, tạo PowerPoint, viết code tối ưu
"""

import os
import uuid
import base64
import tempfile
from datetime import datetime
from typing import Optional
from pathlib import Path

import anthropic
from fastapi import FastAPI, HTTPException, BackgroundTasks
from fastapi.responses import FileResponse, JSONResponse, HTMLResponse
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import json

# ───────────────────────────────────────────
# Khởi tạo
# ───────────────────────────────────────────
app = FastAPI(
    title="🤖 AI API",
    description="API AI thông minh: Chat, tạo PowerPoint, viết Code tối ưu",
    version="2.0.0"
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

client = anthropic.Anthropic(api_key=os.environ.get("ANTHROPIC_API_KEY"))

# Lưu lịch sử hội thoại theo user_id
sessions: dict[str, list] = {}

# ───────────────────────────────────────────
# Models
# ───────────────────────────────────────────
class ChatRequest(BaseModel):
    user_id: str = "default"
    message: str
    system_prompt: Optional[str] = None
    reset: bool = False

class CodeRequest(BaseModel):
    user_id: str = "default"
    task: str
    language: str = "python"
    optimize_level: str = "maximum"  # basic | good | maximum

class PowerPointRequest(BaseModel):
    title: str
    topic: str
    slides: Optional[list[dict]] = None  # [{"title": "...", "content": "..."}]
    num_slides: int = 5
    style: str = "professional"  # professional | modern | minimal | colorful

# ───────────────────────────────────────────
# System Prompts
# ───────────────────────────────────────────
BASE_SYSTEM = """Bạn là một AI thông minh, thân thiện và cực kỳ có ích.

Nguyên tắc trả lời:
- Trả lời ngắn gọn, rõ ràng — không rườm rà
- Làm ĐÚNG theo lệnh người dùng, không suy diễn sai
- Đa dạng cách diễn đạt — không lặp khuôn mẫu
- Thân thiện như người bạn, không cứng nhắc như robot
- Nếu không biết → nói thẳng, không bịa đặt"""

CODE_SYSTEM = """Bạn là senior engineer viết code cực kỳ tối ưu và sạch.

Tiêu chuẩn code:
- Đúng logic, không có bug
- Tối ưu thuật toán (Big-O tốt nhất)
- Clean code: đặt tên rõ ràng, có comment ngắn gọn
- Xử lý edge case và lỗi đầy đủ  
- Kèm ví dụ sử dụng và giải thích ngắn
- Viết theo best practice của ngôn ngữ đó"""

PPTX_SYSTEM = """Bạn tạo nội dung slide PowerPoint chuyên nghiệp.
Trả về JSON THUẦN TÚY (không có ```json), format:
{
  "slides": [
    {
      "title": "Tiêu đề slide",
      "points": ["Điểm 1", "Điểm 2", "Điểm 3"],
      "note": "Ghi chú thuyết trình (tuỳ chọn)"
    }
  ]
}
Nội dung súc tích, chuyên nghiệp, mỗi slide 3-5 điểm."""

# ───────────────────────────────────────────
# Helper: Gọi Claude
# ───────────────────────────────────────────
def call_claude(messages: list, system: str, max_tokens: int = 2048) -> str:
    response = client.messages.create(
        model="claude-sonnet-4-20250514",
        max_tokens=max_tokens,
        system=system,
        messages=messages
    )
    return response.content[0].text

# ───────────────────────────────────────────
# Helper: Tạo file PowerPoint đẹp
# ───────────────────────────────────────────
STYLES = {
    "professional": {
        "bg": RGBColor(0x1A, 0x1A, 0x2E),
        "title_bg": RGBColor(0x16, 0x21, 0x3E),
        "accent": RGBColor(0x0F, 0x3A, 0x5E),
        "title_color": RGBColor(0xE0, 0xF7, 0xFF),
        "text_color": RGBColor(0xCC, 0xE5, 0xFF),
        "bullet_color": RGBColor(0x4F, 0xC3, 0xF7),
        "font_title": "Calibri",
        "font_body": "Calibri Light",
    },
    "modern": {
        "bg": RGBColor(0xF8, 0xF9, 0xFA),
        "title_bg": RGBColor(0x21, 0x29, 0x5C),
        "accent": RGBColor(0xF9, 0x61, 0x67),
        "title_color": RGBColor(0xFF, 0xFF, 0xFF),
        "text_color": RGBColor(0x2D, 0x3A, 0x4A),
        "bullet_color": RGBColor(0xF9, 0x61, 0x67),
        "font_title": "Trebuchet MS",
        "font_body": "Calibri",
    },
    "minimal": {
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "title_bg": RGBColor(0x36, 0x45, 0x4F),
        "accent": RGBColor(0x02, 0xC3, 0x9A),
        "title_color": RGBColor(0xFF, 0xFF, 0xFF),
        "text_color": RGBColor(0x36, 0x45, 0x4F),
        "bullet_color": RGBColor(0x02, 0xC3, 0x9A),
        "font_title": "Cambria",
        "font_body": "Calibri Light",
    },
    "colorful": {
        "bg": RGBColor(0xFF, 0xF8, 0xF0),
        "title_bg": RGBColor(0xB8, 0x50, 0x42),
        "accent": RGBColor(0x02, 0x80, 0x90),
        "title_color": RGBColor(0xFF, 0xFF, 0xFF),
        "text_color": RGBColor(0x2C, 0x1A, 0x1A),
        "bullet_color": RGBColor(0xB8, 0x50, 0x42),
        "font_title": "Georgia",
        "font_body": "Calibri",
    },
}

def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, text, left, top, width, height,
                 font_name, font_size, color, bold=False, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    return txBox

def create_pptx(title: str, slides_data: list, style: str = "professional") -> str:
    s = STYLES.get(style, STYLES["professional"])
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    W = prs.slide_width
    H = prs.slide_height

    # ── Slide tiêu đề ──
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(title_slide, s["bg"])

    # Thanh màu trái
    bar = title_slide.shapes.add_shape(1, 0, 0, Inches(0.5), H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = s["accent"]
    bar.line.fill.background()

    # Tiêu đề lớn
    add_text_box(title_slide, title,
                 Inches(1), Inches(2.5), Inches(10), Inches(1.8),
                 s["font_title"], 44, s["title_color"], bold=True, align=PP_ALIGN.CENTER)

    # Ngày
    add_text_box(title_slide, datetime.now().strftime("%d/%m/%Y"),
                 Inches(1), Inches(4.5), Inches(10), Inches(0.5),
                 s["font_body"], 14, s["bullet_color"], align=PP_ALIGN.CENTER)

    # ── Slides nội dung ──
    for i, slide_info in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide, s["bg"])

        # Header bar
        header = slide.shapes.add_shape(1, 0, 0, W, Inches(1.2))
        header.fill.solid()
        header.fill.fore_color.rgb = s["title_bg"]
        header.line.fill.background()

        # Số slide
        add_text_box(slide, f"{i+1:02d}",
                     Inches(0.2), Inches(0.1), Inches(0.8), Inches(1.0),
                     s["font_title"], 28, s["bullet_color"], bold=True)

        # Tiêu đề slide
        add_text_box(slide, slide_info.get("title", ""),
                     Inches(1.1), Inches(0.15), Inches(11), Inches(0.9),
                     s["font_title"], 26, s["title_color"], bold=True)

        # Thanh accent dưới header
        sep = slide.shapes.add_shape(1, 0, Inches(1.2), W, Pt(3))
        sep.fill.solid()
        sep.fill.fore_color.rgb = s["accent"]
        sep.line.fill.background()

        # Bullet points
        points = slide_info.get("points", [])
        start_y = Inches(1.5)
        row_h = Inches(0.72)

        for j, point in enumerate(points[:6]):
            y = start_y + j * row_h

            # Bullet dot
            dot = slide.shapes.add_shape(9, Inches(0.6), y + Inches(0.18),
                                         Pt(10), Pt(10))
            dot.fill.solid()
            dot.fill.fore_color.rgb = s["bullet_color"]
            dot.line.fill.background()

            # Text
            add_text_box(slide, point,
                         Inches(1.0), y, Inches(11.5), row_h,
                         s["font_body"], 16, s["text_color"])

        # Ghi chú nhỏ
        note = slide_info.get("note", "")
        if note:
            add_text_box(slide, f"💡 {note}",
                         Inches(0.5), Inches(6.8), Inches(12), Inches(0.5),
                         s["font_body"], 11, s["bullet_color"])

        # Số trang góc phải
        add_text_box(slide, f"{i+2} / {len(slides_data)+1}",
                     Inches(11.8), Inches(7.1), Inches(1.3), Inches(0.3),
                     s["font_body"], 10, s["bullet_color"], align=PP_ALIGN.RIGHT)

    # ── Slide kết ──
    end_slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(end_slide, s["title_bg"])
    add_text_box(end_slide, "Cảm ơn!", Inches(1), Inches(3.0), Inches(11), Inches(1.2),
                 s["font_title"], 48, s["title_color"], bold=True, align=PP_ALIGN.CENTER)
    add_text_box(end_slide, title, Inches(1), Inches(4.2), Inches(11), Inches(0.6),
                 s["font_body"], 18, s["bullet_color"], align=PP_ALIGN.CENTER)

    # Lưu file
    output_dir = Path(tempfile.gettempdir()) / "ai_pptx"
    output_dir.mkdir(exist_ok=True)
    filename = output_dir / f"{uuid.uuid4().hex[:8]}_{title[:20].replace(' ','_')}.pptx"
    prs.save(str(filename))
    return str(filename)

# ───────────────────────────────────────────
# Routes
# ───────────────────────────────────────────

@app.get("/", response_class=HTMLResponse)
def chat_ui():
    """Giao diện chat"""
    html_file = Path(__file__).parent / "index.html"
    if html_file.exists():
        return HTMLResponse(content=html_file.read_text(encoding="utf-8"))
    return HTMLResponse("<h1>index.html not found</h1>", status_code=404)

@app.get("/api")
def root():
    return {
        "message": "🤖 AI API đang hoạt động!",
        "endpoints": {
            "POST /chat": "Chat thông minh (đa người dùng)",
            "POST /code": "Viết code tối ưu",
            "POST /powerpoint": "Tạo file PowerPoint",
            "DELETE /session/{user_id}": "Xoá lịch sử hội thoại",
            "GET /sessions": "Xem danh sách sessions"
        }
    }

@app.post("/chat")
async def chat(req: ChatRequest):
    """Chat thông minh — nhớ ngữ cảnh theo từng user_id"""
    if req.reset or req.user_id not in sessions:
        sessions[req.user_id] = []

    sessions[req.user_id].append({"role": "user", "content": req.message})

    system = req.system_prompt or BASE_SYSTEM
    reply = call_claude(sessions[req.user_id], system)

    sessions[req.user_id].append({"role": "assistant", "content": reply})

    return {
        "user_id": req.user_id,
        "reply": reply,
        "turns": len(sessions[req.user_id]) // 2
    }

@app.post("/code")
async def generate_code(req: CodeRequest):
    """Viết code tối ưu theo yêu cầu"""
    system = CODE_SYSTEM
    if req.optimize_level == "maximum":
        system += "\nƯu tiên: tốc độ tối đa, memory tối thiểu, code đẹp nhất có thể."

    if req.user_id not in sessions:
        sessions[req.user_id] = []

    prompt = f"Ngôn ngữ: {req.language}\nYêu cầu: {req.task}"
    sessions[req.user_id].append({"role": "user", "content": prompt})

    code_reply = call_claude(sessions[req.user_id], system, max_tokens=4096)
    sessions[req.user_id].append({"role": "assistant", "content": code_reply})

    return {
        "user_id": req.user_id,
        "language": req.language,
        "optimize_level": req.optimize_level,
        "code": code_reply
    }

@app.post("/powerpoint")
async def create_powerpoint(req: PowerPointRequest):
    """Tạo file PowerPoint đẹp, tải về ngay"""
    # Nếu chưa có slides_data → dùng AI tạo nội dung
    if not req.slides:
        prompt = (
            f"Tạo {req.num_slides} slide cho bài thuyết trình:\n"
            f"Tiêu đề: {req.title}\n"
            f"Chủ đề: {req.topic}\n"
            f"Phong cách: {req.style}\n"
            f"Mỗi slide có 3-5 điểm ngắn gọn, súc tích."
        )
        raw = call_claude([{"role": "user", "content": prompt}], PPTX_SYSTEM, max_tokens=3000)

        try:
            # Claude đôi khi bọc trong ```json ... ```
            clean = raw.strip()
            if clean.startswith("```"):
                clean = clean.split("\n", 1)[1].rsplit("```", 1)[0].strip()
            data = json.loads(clean)
            slides_data = data["slides"]
        except Exception:
            raise HTTPException(status_code=500, detail=f"Lỗi parse AI response: {raw[:200]}")
    else:
        slides_data = req.slides

    filepath = create_pptx(req.title, slides_data, req.style)

    return FileResponse(
        path=filepath,
        filename=f"{req.title.replace(' ', '_')}.pptx",
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

@app.delete("/session/{user_id}")
def reset_session(user_id: str):
    """Xoá lịch sử hội thoại của một user"""
    if user_id in sessions:
        del sessions[user_id]
        return {"message": f"Đã xoá session của {user_id}"}
    return {"message": "Session không tồn tại"}

@app.get("/sessions")
def list_sessions():
    """Xem tất cả sessions đang hoạt động"""
    return {
        uid: {"turns": len(msgs) // 2, "messages": len(msgs)}
        for uid, msgs in sessions.items()
    }
